#!/usr/bin/env python3
"""
Styled PPT Generator
Generates a PPTX file from a JSON content file and local images.
Supports multiple visual styles via layout configurations.
"""

import json
import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert HEX color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_text(slide, text, x, y, w, h, font_size=18, bold=False, color=(0,0,0), align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.alignment = align
    return textbox

def create_slide(prs, slide_data, theme_config):
    """Creates a single slide based on data and theme."""
    layout_idx = 6 # Blank layout usually
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # 1. Background Image (if present)
    image_path = slide_data.get("image_path")
    if image_path and os.path.exists(image_path):
        # If it's a cover or full visuals, use full background
        if slide_data.get("type") == "cover":
            slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        else:
            # For content slides, maybe put image on the right or specific area
            # But recent high-end styles (Lego, Clay) often look good with full bleed visuals
            # Let's try a split layout for content: Left Text, Right Image
            left = prs.slide_width / 2
            slide.shapes.add_picture(image_path, left, 0, prs.slide_width/2, prs.slide_height)

    # 2. Text Content
    title_color = hex_to_rgb(theme_config.get("title_color", "#000000"))
    body_color = hex_to_rgb(theme_config.get("body_color", "#333333"))

    if slide_data.get("type") == "cover":
        # Cover Text (Overlay needed?)
        # For generated full-bleed covers, maybe we don't need text if the image HAS text?
        # But usually we add text. Let's add text with a shadow or box.
        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        
        # Add a semi-transparent box for readability
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.fill.transparency = 0.2
        shape.line.fill.background() # No line

        add_text(slide, title, 1.2, 4.2, 7.6, 1, 44, True, title_color)
        add_text(slide, subtitle, 1.2, 5.2, 7.6, 1, 24, False, body_color)

    else:
        # Content Slide
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        
        # Title
        add_text(slide, title, 0.5, 0.5, 5.5, 1, 32, True, title_color)
        
        # Bullets
        y_pos = 1.8
        for bullet in bullets:
            add_text(slide, "• " + bullet, 0.5, y_pos, 5.5, 0.8, 18, False, body_color)
            y_pos += 0.8

    # 3. Speaker Notes
    notes = slide_data.get("notes", "")
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

def generate_ppt(json_path, output_path):
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme_config = data.get("theme_config", {
        "title_color": "#000000",
        "body_color": "#333333"
    })

    for slide_data in data.get("slides", []):
        create_slide(prs, slide_data, theme_config)

    prs.save(output_path)
    print(f"Generated PPT at: {output_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True, help="Path to JSON file with slide content")
    parser.add_argument("--output", required=True, help="Path to output PPTX file")
    args = parser.parse_args()
    
    generate_ppt(args.input, args.output)
